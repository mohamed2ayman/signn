"""Multi-format text extraction service.

Extracts raw text from uploaded contract documents.  Supports PDF, DOCX,
DOC, XLSX, XLS, PPTX, PPT and plain-text files.
"""

from __future__ import annotations

import io
import os
import subprocess
from typing import Any

from PyPDF2 import PdfReader
from docx import Document as DocxDocument


class TextExtractorService:
    """Extracts text from various file formats."""

    def __init__(self) -> None:
        self.last_page_count: int = 0

    def extract(self, file_path: str, mime_type: str) -> dict[str, Any]:
        """Extract text from the file at *file_path*.

        Returns
        -------
        dict with keys ``text`` (str) and ``page_count`` (int).
        """
        mime = mime_type.lower()
        self.last_page_count = 0

        if mime in ("application/pdf",):
            return self._extract_pdf(file_path)
        elif mime in (
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "application/docx",
        ):
            return self._extract_docx(file_path)
        elif mime in ("application/msword",):
            return self._extract_doc(file_path)
        elif mime in (
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "application/vnd.ms-excel",
        ):
            return self._extract_xlsx(file_path)
        elif mime in (
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "application/vnd.ms-powerpoint",
        ):
            return self._extract_pptx(file_path)
        elif mime.startswith("text/"):
            return self._extract_txt(file_path)
        else:
            # Fallback: try reading as text
            return self._extract_txt(file_path)

    # ------------------------------------------------------------------
    # PDF
    # ------------------------------------------------------------------

    def _extract_pdf(self, path: str) -> dict[str, Any]:
        reader = PdfReader(path)
        self.last_page_count = len(reader.pages)
        pages_text: list[str] = []

        for page in reader.pages:
            text = page.extract_text() or ""
            pages_text.append(text)

        full_text = "\n\n".join(pages_text)

        # If no text was extracted (scanned PDF), attempt OCR
        if not full_text.strip():
            full_text = self._ocr_pdf(path)

        return {"text": full_text, "page_count": self.last_page_count}

    def _ocr_pdf(self, path: str) -> str:
        """Attempt OCR on a scanned PDF.

        Strategy:
        1. Try ``pdftotext`` (poppler-utils) — fast, works for PDFs with
           embedded text layers that PyPDF2 failed to read.
        2. Fall back to ``pytesseract`` + ``pdf2image`` — renders each page
           as an image and runs Tesseract OCR with Arabic + English support.
        """

        # --- Attempt 1: pdftotext (poppler) ---
        try:
            result = subprocess.run(
                ["pdftotext", path, "-"],
                capture_output=True,
                text=True,
                timeout=120,
            )
            if result.returncode == 0 and result.stdout.strip():
                return result.stdout
        except (FileNotFoundError, subprocess.TimeoutExpired):
            pass

        # --- Attempt 2: pytesseract + pdf2image (chunked) ---
        try:
            import gc
            import logging
            import pytesseract
            from pdf2image import convert_from_path, pdfinfo_from_path

            logger = logging.getLogger(__name__)

            # Get total page count without loading images
            try:
                info = pdfinfo_from_path(path)
                total_pages = info.get("Pages", 0)
            except Exception:
                total_pages = self.last_page_count or 50

            pages_text: list[str] = []
            batch_size = 5  # Process 5 pages at a time to limit memory

            for batch_start in range(1, total_pages + 1, batch_size):
                batch_end = min(batch_start + batch_size - 1, total_pages)
                try:
                    images = convert_from_path(
                        path,
                        dpi=150,
                        first_page=batch_start,
                        last_page=batch_end,
                    )
                    for img in images:
                        try:
                            text = pytesseract.image_to_string(img, lang="ara+eng")
                            if text.strip():
                                pages_text.append(text)
                        except Exception as page_exc:
                            logger.warning(
                                "OCR failed for page in batch %d-%d: %s",
                                batch_start, batch_end, page_exc,
                            )
                        finally:
                            img.close()
                    del images
                    gc.collect()
                except Exception as batch_exc:
                    logger.warning(
                        "OCR batch %d-%d failed: %s",
                        batch_start, batch_end, batch_exc,
                    )
                    continue

            if pages_text:
                return "\n\n".join(pages_text)
        except (ImportError, Exception) as exc:
            # Log but don't crash — OCR is best-effort
            import logging
            logging.getLogger(__name__).warning("OCR fallback failed: %s", exc)

        return ""

    # ------------------------------------------------------------------
    # DOCX
    # ------------------------------------------------------------------

    def _extract_docx(self, path: str) -> dict[str, Any]:
        doc = DocxDocument(path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]

        # Extract text from table cells (contract terms, payment schedules, etc.
        # are often in tables and are missed by doc.paragraphs alone)
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        paragraphs.append(cell.text.strip())

        # Extract headers (document title, article headings in header band)
        for section in doc.sections:
            header = section.header
            for para in header.paragraphs:
                if para.text.strip():
                    paragraphs.append(para.text.strip())

        self.last_page_count = max(1, len(paragraphs) // 40)  # Estimate
        return {"text": "\n\n".join(paragraphs), "page_count": self.last_page_count}

    # ------------------------------------------------------------------
    # DOC (legacy Word)
    # ------------------------------------------------------------------

    def _extract_doc(self, path: str) -> dict[str, Any]:
        # Try macOS textutil first
        try:
            result = subprocess.run(
                ["textutil", "-convert", "txt", "-stdout", path],
                capture_output=True,
                text=True,
                timeout=30,
            )
            if result.returncode == 0:
                text = result.stdout
                self.last_page_count = max(1, len(text.split("\n")) // 50)
                return {"text": text, "page_count": self.last_page_count}
        except (FileNotFoundError, subprocess.TimeoutExpired):
            pass

        # Fallback: try reading as docx (some .doc files are actually docx)
        try:
            return self._extract_docx(path)
        except Exception:
            return {"text": "", "page_count": 0}

    # ------------------------------------------------------------------
    # XLSX / XLS
    # ------------------------------------------------------------------

    def _extract_xlsx(self, path: str) -> dict[str, Any]:
        from openpyxl import load_workbook

        wb = load_workbook(path, read_only=True, data_only=True)
        sheets_text: list[str] = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows_text: list[str] = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                line = " | ".join(cells).strip()
                if line and line != " | ".join([""] * len(cells)):
                    rows_text.append(line)
            if rows_text:
                sheets_text.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows_text))

        wb.close()
        full_text = "\n\n".join(sheets_text)
        self.last_page_count = len(wb.sheetnames)
        return {"text": full_text, "page_count": self.last_page_count}

    # ------------------------------------------------------------------
    # PPTX / PPT
    # ------------------------------------------------------------------

    def _extract_pptx(self, path: str) -> dict[str, Any]:
        from pptx import Presentation

        prs = Presentation(path)
        slides_text: list[str] = []

        for i, slide in enumerate(prs.slides, 1):
            texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            texts.append(text)
            if texts:
                slides_text.append(f"[Slide {i}]\n" + "\n".join(texts))

        self.last_page_count = len(prs.slides)
        return {"text": "\n\n".join(slides_text), "page_count": self.last_page_count}

    # ------------------------------------------------------------------
    # Plain text
    # ------------------------------------------------------------------

    def _extract_txt(self, path: str) -> dict[str, Any]:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            text = f.read()
        self.last_page_count = max(1, len(text.split("\n")) // 50)
        return {"text": text, "page_count": self.last_page_count}
