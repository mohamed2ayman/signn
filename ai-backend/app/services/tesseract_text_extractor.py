"""Tesseract-based text extraction (default backend).

Implements the full three-strategy PDF extraction cascade:
  1. PyPDF2  — fast; works for PDFs with a native text layer.
  2. pdftotext (poppler-utils subprocess) — fast; works for PDFs that have
     an embedded text layer that PyPDF2 failed to parse.
  3. pytesseract + pdf2image — full OCR for scanned PDFs, with Arabic + English
     language packs enabled (``ara+eng``).

Non-PDF formats (DOCX, DOC, XLSX, PPTX, plain text) are also handled by this
class — they are not part of the :class:`BaseTextExtractor` interface since they
do not vary across extraction backends.
"""

from __future__ import annotations

import logging
import re
import subprocess
from typing import Any

from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from docx.oxml.ns import qn

from app.config.settings import get_settings
from app.services.base_text_extractor import BaseTextExtractor

_logger = logging.getLogger(__name__)

# DPI for the force-OCR path.  300 produced clean Arabic in the Phase 7.27
# investigation; lower DPI tested poorly.  Single edit point for future tuning.
FORCE_OCR_DPI = 300


# ---------------------------------------------------------------------------
# Word list-numbering reconstruction (Issue 1 — Defects A + C)
# ---------------------------------------------------------------------------
# python-docx's ``paragraph.text`` returns ONLY the concatenated run text. A
# Word list item's displayed marker — the auto-number ("1.", "1)", "1-1") OR the
# bullet glyph ("•") — lives in the numbering definitions (numbering.xml) via the
# paragraph's ``numPr`` reference, NOT in the run text. So auto-numbered clauses
# arrive number-less (→ empty section_number, Defect C) and bulleted lists arrive
# structure-less (→ flattening, Defect A). The helper below rebuilds each marker
# from the <w:num>/<w:abstractNum>/<w:lvl> definitions plus a running per-list
# counter, so the model SEES the numbering/bullets and can preserve them.
#
# Best-effort by design: ANY resolution failure falls back to "no label" (the
# pre-fix behaviour). Text extraction must never crash on unusual numbering XML.


def _paragraph_numpr(paragraph: Any) -> tuple[str, int] | None:
    """Return ``(numId, ilvl)`` when *paragraph* is a Word list item, else None.

    Pure XML navigation (``w:pPr/w:numPr``) so it is robust across python-docx
    versions and never raises on a missing element.
    """
    try:
        p_el = paragraph._p
        pPr = p_el.find(qn("w:pPr"))
        if pPr is None:
            return None
        numPr = pPr.find(qn("w:numPr"))
        if numPr is None:
            return None
        num_id_el = numPr.find(qn("w:numId"))
        if num_id_el is None:
            return None
        num_id = num_id_el.get(qn("w:val"))
        if not num_id:
            return None
        ilvl_el = numPr.find(qn("w:ilvl"))
        ilvl_val = ilvl_el.get(qn("w:val")) if ilvl_el is not None else None
        ilvl = int(ilvl_val) if ilvl_val is not None else 0
        return num_id, ilvl
    except Exception:  # noqa: BLE001 — never crash extraction on odd XML
        return None


class _DocxListNumbering:
    """Reconstructs the rendered marker for Word list items (see module note).

    Parses the document's numbering part once, then renders labels on demand
    while advancing a running per-``(numId, level)`` counter — mirroring how Word
    displays the list. Numbered levels honour their ``numFmt`` (decimal / letter /
    roman) and ``lvlText`` template (e.g. ``"%1."`` → ``"1."``, ``"%1-%2"`` →
    ``"1-1"``); bullet levels render ``"• "``.
    """

    _ROMAN = [
        (1000, "m"), (900, "cm"), (500, "d"), (400, "cd"), (100, "c"),
        (90, "xc"), (50, "l"), (40, "xl"), (10, "x"), (9, "ix"),
        (5, "v"), (4, "iv"), (1, "i"),
    ]

    def __init__(self, doc: Any) -> None:
        self._num_to_abs: dict[str, str] = {}                 # numId -> abstractNumId
        self._abs_levels: dict[str, dict[int, dict[str, Any]]] = {}  # absId -> {ilvl: lvl}
        self._counters: dict[str, dict[int, int]] = {}        # numId -> {ilvl: current}
        try:
            self._load(doc)
        except Exception:  # noqa: BLE001 — a doc with no/odd numbering just gets no labels
            self._num_to_abs = {}
            self._abs_levels = {}

    def _load(self, doc: Any) -> None:
        numbering = doc.part.numbering_part.element  # raises if the doc has no lists
        for num in numbering.findall(qn("w:num")):
            num_id = num.get(qn("w:numId"))
            abs_ref = num.find(qn("w:abstractNumId"))
            if num_id is not None and abs_ref is not None:
                self._num_to_abs[num_id] = abs_ref.get(qn("w:val"))
        for abs_num in numbering.findall(qn("w:abstractNum")):
            abs_id = abs_num.get(qn("w:abstractNumId"))
            levels: dict[int, dict[str, Any]] = {}
            for lvl in abs_num.findall(qn("w:lvl")):
                try:
                    ilvl = int(lvl.get(qn("w:ilvl")) or 0)
                except (TypeError, ValueError):
                    ilvl = 0
                fmt_el = lvl.find(qn("w:numFmt"))
                txt_el = lvl.find(qn("w:lvlText"))
                start_el = lvl.find(qn("w:start"))
                start_val = start_el.get(qn("w:val")) if start_el is not None else None
                levels[ilvl] = {
                    "numFmt": fmt_el.get(qn("w:val")) if fmt_el is not None else "decimal",
                    "lvlText": txt_el.get(qn("w:val")) if txt_el is not None else "%1.",
                    "start": int(start_val) if start_val and start_val.isdigit() else 1,
                }
            self._abs_levels[abs_id] = levels

    @property
    def available(self) -> bool:
        """True when the document carries resolvable numbering definitions."""
        return bool(self._abs_levels)

    def _fmt_num(self, n: int, fmt: str) -> str:
        if n < 1:
            n = 1
        if fmt in ("lowerLetter", "upperLetter"):
            s, x = "", n
            while x > 0:
                x, r = divmod(x - 1, 26)
                s = chr(ord("a") + r) + s
            return s.upper() if fmt == "upperLetter" else s
        if fmt in ("lowerRoman", "upperRoman"):
            x, out = n, ""
            for val, sym in self._ROMAN:
                while x >= val:
                    out += sym
                    x -= val
            return out.upper() if fmt == "upperRoman" else out
        # decimal / arabicAlpha / anything else → plain integer
        return str(n)

    def label(self, num_id: str, ilvl: int) -> str | None:
        """Rendered marker (e.g. ``"1. "``, ``"1-1 "``, ``"• "``) for the NEXT
        item of list *num_id* at level *ilvl*, advancing the counters. Returns
        None when this list/level can't be resolved (caller then adds no label).
        """
        abs_id = self._num_to_abs.get(num_id)
        if abs_id is None:
            return None
        levels = self._abs_levels.get(abs_id)
        if not levels or ilvl not in levels:
            return None
        lvl = levels[ilvl]

        # Advance THIS level's counter; reset all DEEPER levels (Word restarts a
        # nested list each time its parent advances). Ancestor levels are kept.
        ctr = self._counters.setdefault(num_id, {})
        ctr[ilvl] = ctr.get(ilvl, lvl["start"] - 1) + 1
        for deeper in [d for d in ctr if d > ilvl]:
            del ctr[deeper]

        if lvl["numFmt"] == "bullet":
            return "• "

        # Numbered: substitute %1..%9 in lvlText with each level's counter,
        # each formatted per that level's own numFmt (so "%1.%2" → "1-b" etc.).
        def _repl(m: "re.Match[str]") -> str:
            lv = int(m.group(1)) - 1  # %1 → level 0
            lv_def = levels.get(lv, {})
            n = ctr.get(lv, lv_def.get("start", 1))
            return self._fmt_num(int(n), lv_def.get("numFmt", "decimal"))

        rendered = re.sub(r"%(\d)", _repl, lvl["lvlText"] or "%1.").strip()
        return f"{rendered} " if rendered else None


class TesseractTextExtractor(BaseTextExtractor):
    """Text extractor using Tesseract OCR as the PDF fallback.

    A new instance is created per Celery task invocation.  This class
    intentionally carries *no* shared mutable state between concurrent tasks
    — the old ``self.last_page_count`` side-effect has been eliminated by
    passing ``page_count`` as an explicit parameter to :meth:`_ocr_pdf`.
    """

    # ------------------------------------------------------------------
    # Public entry point
    # ------------------------------------------------------------------

    def extract(
        self,
        file_path: str,
        mime_type: str,
        force_ocr: bool = False,
    ) -> dict[str, Any]:
        """Extract text from *file_path* according to *mime_type*.

        Parameters
        ----------
        force_ocr:
            When True and the file is a PDF, bypass the digital text-layer fast
            path entirely and OCR rendered page images @ ``FORCE_OCR_DPI``.
            Used for sources whose embedded fonts have a broken ToUnicode CMap
            (e.g. ETA's kaf→آ corruption), where the text layer is unusable but
            non-empty.  Ignored for non-PDF formats.

        Returns
        -------
        dict with keys:
          - ``text`` (str)
          - ``page_count`` (int)
          - ``quality_flags`` (list[str]) — scan quality signals; always empty
            for non-PDF formats (no OCR involved).  Contains ``'ocr_forced'``
            when the force-OCR path was taken.
        """
        mime = mime_type.lower()

        if mime == "application/pdf":
            if force_ocr:
                return self._extract_pdf_force_ocr(file_path)
            return self._extract_pdf(file_path)

        # Non-PDF formats go through dedicated extractors.  No OCR is involved
        # so quality_flags is always empty — add the key for a uniform response
        # shape so callers never need to branch on mime type.
        if mime in (
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "application/docx",
        ):
            result = self._extract_docx(file_path)
        elif mime == "application/msword":
            result = self._extract_doc(file_path)
        elif mime in (
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "application/vnd.ms-excel",
        ):
            result = self._extract_xlsx(file_path)
        elif mime in (
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "application/vnd.ms-powerpoint",
        ):
            result = self._extract_pptx(file_path)
        elif mime.startswith("text/"):
            result = self._extract_txt(file_path)
        else:
            # Fallback: try reading as plain text
            result = self._extract_txt(file_path)

        result.setdefault("quality_flags", [])
        return result

    # ------------------------------------------------------------------
    # PDF — private routing helper
    # ------------------------------------------------------------------

    def _extract_pdf(self, path: str) -> dict[str, Any]:
        """Determine page count, then delegate text extraction to extract_pdf().

        Returns a dict with keys ``text``, ``page_count``, and
        ``quality_flags`` (list of flag strings, empty when the scan is clean
        or the file has a native text layer).
        """
        reader = PdfReader(path)
        page_count = len(reader.pages)
        text, quality_flags = self.extract_pdf(path, page_count)
        return {"text": text, "page_count": page_count, "quality_flags": quality_flags}

    # ------------------------------------------------------------------
    # PDF — forced OCR path (broken-text-layer sources)
    # ------------------------------------------------------------------

    def _extract_pdf_force_ocr(self, path: str) -> dict[str, Any]:
        """Render and OCR the PDF **one page at a time** @ FORCE_OCR_DPI (Arabic).

        Bypasses the PDF text layer entirely — the correct strategy for PDFs
        whose embedded fonts carry a broken ToUnicode CMap (the text layer is
        present but corrupt, so the normal fast path would silently return
        garbage).  Tesseract reads the rendered pixels and produces clean,
        logical-order Arabic.

        MEMORY: pages are rendered individually via pdf2image's
        ``first_page``/``last_page`` so peak RAM stays flat (~30 MB/iteration)
        regardless of page count.  Rendering all pages at once previously held
        ~2.6 GB for a 100-page doc and OOM-killed the 3 GB worker.

        RESILIENCE: a single page's OCR failure is logged, flagged
        (``ocr_page_<n>_failed``) and skipped (empty text for that page) rather
        than failing the whole document — one missing page is far better than a
        lost corpus document.

        Returns the same dict shape as the text-layer path, with an extra
        ``'ocr_forced'`` quality flag so downstream code can see the path taken.
        """
        from pdf2image import convert_from_path, pdfinfo_from_path
        import pytesseract

        # Page count via pdfinfo_from_path (ships with pdf2image; uses poppler's
        # pdfinfo — no new dependency, no full render just to count).
        total_pages = int(pdfinfo_from_path(path)["Pages"])

        pages_text: list[str] = []
        quality_flags: list[str] = ["ocr_forced"]

        for page_num in range(1, total_pages + 1):
            if page_num % 10 == 1 or page_num == total_pages:
                _logger.info("OCR page %d/%d", page_num, total_pages)
            try:
                images = convert_from_path(
                    path,
                    dpi=FORCE_OCR_DPI,
                    first_page=page_num,
                    last_page=page_num,
                )
                try:
                    pages_text.append(
                        pytesseract.image_to_string(images[0], lang="ara")
                    )
                finally:
                    for img in images:
                        img.close()
            except Exception as exc:  # noqa: BLE001 — one bad page must not kill the doc
                _logger.warning("OCR failed on page %d: %s", page_num, exc)
                quality_flags.append(f"ocr_page_{page_num}_failed")
                pages_text.append("")  # keep page alignment; gap, not failure

        text = "\n\n".join(pages_text)
        return {
            "text": text,
            "page_count": total_pages,
            "quality_flags": quality_flags,
        }

    # ------------------------------------------------------------------
    # PDF — BaseTextExtractor implementation
    # ------------------------------------------------------------------

    def extract_pdf(self, file_path: str, page_count: int) -> tuple[str, list[str]]:
        """Three-strategy PDF text extraction: PyPDF2 → pdftotext → pytesseract.

        Parameters
        ----------
        file_path:
            Path to the PDF file.
        page_count:
            Page count determined by the caller (from a prior PdfReader open).
            Used by :meth:`_ocr_pdf` as a fallback when ``pdfinfo_from_path``
            fails — eliminates the old ``self.last_page_count`` side-effect.

        Returns
        -------
        (text, quality_flags) — quality_flags is empty when the file has a
        native text layer (no OCR was needed).
        """
        reader = PdfReader(file_path)
        pages_text: list[str] = []
        for page in reader.pages:
            pages_text.append(page.extract_text() or "")
        full_text = "\n\n".join(pages_text)

        # If no text was extracted (scanned PDF), attempt OCR
        if not full_text.strip():
            # Prefer the count from the just-opened reader; fall back to the
            # caller-supplied hint if PyPDF2 returned 0 for some reason.
            actual_count = len(reader.pages) or page_count
            return self._ocr_pdf(file_path, actual_count)

        # Digital text layer: no OCR, no quality concerns.
        return full_text, []

    # ------------------------------------------------------------------
    # Phase 7.25 — Scan quality detection and enhancement
    # ------------------------------------------------------------------

    def _assess_quality(self, images: list) -> list[str]:
        """Assess scan quality of a batch of PIL images.

        Samples the first 2 pages (or all pages when fewer than 2).
        Returns a list of quality flag strings.  Each flag encodes both
        the signal type and its measured value, e.g. ``"blur:32.1"``.
        An empty list means the scan passed all quality checks.

        Metrics (using only Pillow + numpy — no opencv required):
        - Blur   : Laplacian variance on grayscale pixel array (pure numpy).
                   Low variance == blurry. Default threshold: 50.
        - Contrast: PIL ImageStat stddev on grayscale channel.
                   Low stddev == washed-out. Default threshold: 20.
        - Skew   : pytesseract.image_to_osd on page 1 only.
                   Rotation >= 10° == skewed. OSD failures are silently
                   swallowed — never block extraction.
        """
        try:
            import numpy as np
            from PIL import ImageStat
        except ImportError:
            # numpy / Pillow absent — skip quality check gracefully
            return []

        settings = get_settings()
        flags: list[str] = []

        sample = images[:2]  # first 2 pages at most

        for idx, image in enumerate(sample):
            # ── Blur check ──────────────────────────────────────────────
            try:
                gray = np.array(image.convert("L"), dtype=np.float32)
                # Pure-numpy Laplacian via convolution kernel (no opencv).
                kernel = np.array([[0, 1, 0], [1, -4, 1], [0, 1, 0]], dtype=np.float32)
                # Manual 2-D convolution using stride tricks (valid region).
                h, w = gray.shape
                kh, kw = kernel.shape
                out_h, out_w = h - kh + 1, w - kw + 1
                shape = (out_h, out_w, kh, kw)
                strides = (gray.strides[0], gray.strides[1], gray.strides[0], gray.strides[1])
                patches = np.lib.stride_tricks.as_strided(gray, shape=shape, strides=strides)
                laplacian = (patches * kernel).sum(axis=(2, 3))
                variance = float(laplacian.var())
                if variance < settings.BLUR_THRESHOLD:
                    flags.append(f"blur:{variance:.1f}")
            except Exception as exc:
                _logger.debug("Blur check failed on page %d: %s", idx, exc)

            # ── Contrast check ──────────────────────────────────────────
            try:
                stat = ImageStat.Stat(image.convert("L"))
                stddev = float(stat.stddev[0])
                if stddev < settings.CONTRAST_THRESHOLD:
                    flags.append(f"contrast:{stddev:.1f}")
            except Exception as exc:
                _logger.debug("Contrast check failed on page %d: %s", idx, exc)

            # ── Skew / rotation check (page 1 only, OSD lang pack required) ──
            if idx == 0:
                try:
                    import pytesseract
                    osd = pytesseract.image_to_osd(
                        image,
                        lang="osd",
                        output_type=pytesseract.Output.DICT,
                    )
                    rotation = abs(float(osd.get("rotate", 0)))
                    if rotation >= settings.ROTATION_THRESHOLD:
                        flags.append(f"rotation:{rotation:.0f}")
                except Exception as exc:
                    # OSD can fail when the page has no text at all or the
                    # OSD language pack is missing — skip, never block.
                    _logger.debug("OSD skew check failed: %s", exc)

        return flags

    def _enhance_image(self, image, flags: list[str]):
        """Apply opportunistic enhancements based on quality flags.

        Only touch the image when a relevant flag is present — avoids
        degrading high-quality inputs.

        - Low contrast  → PIL.ImageOps.autocontrast(cutoff=2)
        - Rotation ≥ 5° → rotate to correct skew using the detected angle
        """
        try:
            from PIL import ImageOps
        except ImportError:
            return image

        enhanced = image

        has_contrast = any(f.startswith("contrast:") for f in flags)
        if has_contrast:
            try:
                enhanced = ImageOps.autocontrast(enhanced.convert("L"), cutoff=2).convert("RGB")
            except Exception as exc:
                _logger.debug("autocontrast failed: %s", exc)

        for flag in flags:
            if flag.startswith("rotation:"):
                try:
                    degrees = float(flag.split(":")[1])
                    if abs(degrees) >= 5:
                        enhanced = enhanced.rotate(
                            -degrees, expand=True, fillcolor=(255, 255, 255)
                        )
                except Exception as exc:
                    _logger.debug("Rotation correction failed: %s", exc)
                break

        return enhanced

    def _ocr_pdf(self, path: str, page_count: int) -> tuple[str, list[str]]:
        """Attempt OCR on a scanned PDF.

        Strategy:
        1. ``pdftotext`` (poppler-utils) — fast; for PDFs with an embedded
           text layer that PyPDF2 failed to read.
        2. ``pytesseract`` + ``pdf2image`` — renders each page as an image
           and runs Tesseract OCR with Arabic + English support.  Quality
           check fires after the first batch is rendered; enhancements are
           applied when flags are detected.

        Parameters
        ----------
        path:
            Path to the PDF file.
        page_count:
            Total page count passed explicitly.  Used as a fallback when
            ``pdfinfo_from_path`` fails — no shared instance state required.

        Returns
        -------
        (text, quality_flags) tuple where quality_flags is a list of flag
        strings (e.g. ``["blur:32.1", "contrast:15.4"]``).  An empty list
        means the scan passed all quality checks.
        """

        # --- Attempt 1: pdftotext (poppler) ---
        try:
            result = subprocess.run(
                ["pdftotext", path, "-"],
                capture_output=True,
                text=True,
                timeout=120,
            )
            if result.returncode == 0 and result.stdout.strip():
                # Digital text layer found — no OCR needed, no quality issues.
                return result.stdout, []
        except (FileNotFoundError, subprocess.TimeoutExpired):
            pass

        # --- Attempt 2: pytesseract + pdf2image (chunked) ---
        # All OCR imports are lazy so the service degrades gracefully when
        # these optional packages are absent (ImportError is caught below).
        try:
            import gc
            import pytesseract
            from pdf2image import convert_from_path, pdfinfo_from_path

            # Get total page count without loading images
            try:
                info = pdfinfo_from_path(path)
                total_pages = info.get("Pages", 0)
            except Exception:
                # Fall back to the page_count passed by the caller
                total_pages = page_count or 50

            pages_text: list[str] = []
            batch_size = 5  # Process 5 pages at a time to limit memory
            quality_flags: list[str] = []
            quality_checked = False

            for batch_start in range(1, total_pages + 1, batch_size):
                batch_end = min(batch_start + batch_size - 1, total_pages)
                try:
                    images = convert_from_path(
                        path,
                        dpi=150,
                        first_page=batch_start,
                        last_page=batch_end,
                    )

                    # Quality check + enhancement on the first batch only.
                    # Cost: < 300 ms for 2-page sample; negligible vs pdf2image.
                    if not quality_checked:
                        quality_flags = self._assess_quality(images)
                        quality_checked = True
                        if quality_flags:
                            _logger.info(
                                "Scan quality flags detected for %s: %s",
                                path,
                                quality_flags,
                            )
                            # Apply enhancements to first-batch images in-place.
                            images = [self._enhance_image(img, quality_flags) for img in images]

                    for img in images:
                        try:
                            text = pytesseract.image_to_string(img, lang="ara+eng")
                            if text.strip():
                                pages_text.append(text)
                        except Exception as page_exc:
                            _logger.warning(
                                "OCR failed for page in batch %d-%d: %s",
                                batch_start,
                                batch_end,
                                page_exc,
                            )
                        finally:
                            img.close()
                    del images
                    gc.collect()
                except Exception as batch_exc:
                    _logger.warning(
                        "OCR batch %d-%d failed: %s",
                        batch_start,
                        batch_end,
                        batch_exc,
                    )
                    continue

            if pages_text:
                return "\n\n".join(pages_text), quality_flags
            return "", quality_flags
        except (ImportError, Exception) as exc:
            # Log but don't crash — OCR is best-effort
            _logger.warning("OCR fallback failed: %s", exc)

        return "", []

    # ------------------------------------------------------------------
    # DOCX
    # ------------------------------------------------------------------

    def _extract_docx(self, path: str) -> dict[str, Any]:
        doc = DocxDocument(path)

        # Issue 1 (Defects A + C): reconstruct Word list markers dropped by
        # ``p.text``. Numbered clauses regain their number (so the model can fill
        # section_number); bulleted lists regain "• " (so structure survives).
        numbering = _DocxListNumbering(doc)
        paragraphs: list[str] = []
        for p in doc.paragraphs:
            text = p.text
            if not text.strip():
                continue
            info = _paragraph_numpr(p)
            if info is not None:
                label = numbering.label(info[0], info[1])
                # Guard against double-prefixing when the marker was ALSO typed
                # literally into the run text (e.g. a manual "1-").
                if label and not text.lstrip().startswith(label.strip()):
                    text = f"{label}{text}"
            paragraphs.append(text)

        # Extract text from table cells (contract terms, payment schedules, etc.
        # are often in tables and are missed by doc.paragraphs alone)
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        paragraphs.append(cell.text.strip())

        # Extract headers (document title, article headings in header band)
        for section in doc.sections:
            header = section.header
            for para in header.paragraphs:
                if para.text.strip():
                    paragraphs.append(para.text.strip())

        page_count = max(1, len(paragraphs) // 40)  # Estimate
        return {"text": "\n\n".join(paragraphs), "page_count": page_count}

    # ------------------------------------------------------------------
    # DOC (legacy Word)
    # ------------------------------------------------------------------

    def _extract_doc(self, path: str) -> dict[str, Any]:
        # Try macOS textutil first
        try:
            result = subprocess.run(
                ["textutil", "-convert", "txt", "-stdout", path],
                capture_output=True,
                text=True,
                timeout=30,
            )
            if result.returncode == 0:
                text = result.stdout
                page_count = max(1, len(text.split("\n")) // 50)
                return {"text": text, "page_count": page_count}
        except (FileNotFoundError, subprocess.TimeoutExpired):
            pass

        # Fallback: try reading as docx (some .doc files are actually docx)
        try:
            return self._extract_docx(path)
        except Exception:
            return {"text": "", "page_count": 0}

    # ------------------------------------------------------------------
    # XLSX / XLS
    # ------------------------------------------------------------------

    def _extract_xlsx(self, path: str) -> dict[str, Any]:
        from openpyxl import load_workbook

        wb = load_workbook(path, read_only=True, data_only=True)
        sheets_text: list[str] = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows_text: list[str] = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                line = " | ".join(cells).strip()
                if line and line != " | ".join([""] * len(cells)):
                    rows_text.append(line)
            if rows_text:
                sheets_text.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows_text))

        wb.close()
        full_text = "\n\n".join(sheets_text)
        page_count = len(wb.sheetnames)
        return {"text": full_text, "page_count": page_count}

    # ------------------------------------------------------------------
    # PPTX / PPT
    # ------------------------------------------------------------------

    def _extract_pptx(self, path: str) -> dict[str, Any]:
        from pptx import Presentation

        prs = Presentation(path)
        slides_text: list[str] = []

        for i, slide in enumerate(prs.slides, 1):
            texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            texts.append(text)
            if texts:
                slides_text.append(f"[Slide {i}]\n" + "\n".join(texts))

        page_count = len(prs.slides)
        return {"text": "\n\n".join(slides_text), "page_count": page_count}

    # ------------------------------------------------------------------
    # Plain text
    # ------------------------------------------------------------------

    def _extract_txt(self, path: str) -> dict[str, Any]:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            text = f.read()
        page_count = max(1, len(text.split("\n")) // 50)
        return {"text": text, "page_count": page_count}
